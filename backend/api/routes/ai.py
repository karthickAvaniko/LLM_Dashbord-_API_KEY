from fastapi import APIRouter, Header, Form, File, UploadFile, HTTPException, Depends
from fastapi.responses import StreamingResponse, JSONResponse
import httpx, base64, json, traceback
from typing import List, Optional
from backend.core.config import settings
from backend.core.schemas import get_mode, list_modes
from backend.schemas.api_key import ChatRequest, GenerateRequest
from backend.api.dependencies import verify_key, log_usage

router = APIRouter()


def _extract_text(choice: dict) -> str:
    msg = choice.get("message", {}) or {}
    text = msg.get("content")
    if not text:
        text = msg.get("reasoning_content") or msg.get("reasoning") or ""
    if not text:
        text = choice.get("text", "") or ""
    return text or ""


def _build_payload(messages, max_tokens, temperature, mode=None, enable_thinking=False):
    """Build vLLM payload. If `mode` matches an active DB mode, apply its schema + system prompt.
    Thinking is DISABLED by default — production wants direct output, not reasoning leak."""
    mode_def = get_mode(mode) if mode else None
    if mode_def:
        sys_msg = mode_def.get("system_prompt") or ""
        if sys_msg and (not messages or messages[0].get("role") != "system"):
            messages = [{"role": "system", "content": sys_msg}] + messages
        max_tokens = max_tokens or mode_def.get("max_tokens") or 8192
        temperature = float(mode_def.get("temperature", 0.0))
    payload = {
        "model": settings.MODEL_NAME,
        "messages": messages,
        "max_tokens": max_tokens,
        "temperature": temperature,
        # Disable Qwen thinking mode — direct response, no reasoning preamble
        "chat_template_kwargs": {"enable_thinking": bool(enable_thinking)},
    }
    if mode_def and mode_def.get("json_schema"):
        payload["guided_json"] = mode_def["json_schema"]
    return payload


# ──────────────────────────────────────────────────────────────────
#  Thinking-tag stripper (belt-and-suspenders defense)
#  Even if vLLM leaks <think>...</think>, we filter before client sees it.
# ──────────────────────────────────────────────────────────────────
class ThinkFilter:
    """Streaming-aware filter that drops everything between <think> and </think>.
    Use one instance per request. Feed deltas in order, get clean output back."""
    def __init__(self):
        self.inside_think = False
        self.buffer = ""

    def feed(self, piece: str) -> str:
        """Feed a delta; returns the cleaned portion (may be empty)."""
        if not piece:
            return ""
        self.buffer += piece
        out = []
        while self.buffer:
            if self.inside_think:
                idx = self.buffer.find("</think>")
                if idx == -1:
                    # Still inside think — wait for more
                    # Keep last 8 chars in case </think> spans deltas
                    if len(self.buffer) > 8:
                        self.buffer = self.buffer[-8:]
                    return "".join(out)
                # Skip past </think>
                self.buffer = self.buffer[idx + len("</think>"):]
                self.inside_think = False
            else:
                idx = self.buffer.find("<think>")
                if idx == -1:
                    # Hold last 7 chars in case <think> spans deltas
                    if len(self.buffer) > 7:
                        out.append(self.buffer[:-7])
                        self.buffer = self.buffer[-7:]
                    return "".join(out)
                # Emit text before <think>, then enter think mode
                out.append(self.buffer[:idx])
                self.buffer = self.buffer[idx + len("<think>"):]
                self.inside_think = True
        return "".join(out)

    def flush(self) -> str:
        """Call at end-of-stream to flush any remaining buffer."""
        if self.inside_think:
            return ""  # incomplete think block — drop it
        out = self.buffer
        self.buffer = ""
        return out


def _strip_think_tags(text: str) -> str:
    """Synchronous version for non-streaming endpoints."""
    f = ThinkFilter()
    out = f.feed(text)
    out += f.flush()
    return out.strip()


# ──────────────────────────────────────────────────────────────────
#  Token estimation + text chunking + Map-Reduce pipeline
# ──────────────────────────────────────────────────────────────────
_CHARS_PER_TOKEN = 3        # conservative avg (Latin ~4, CJK ~1.5)
_CHUNK_TRIGGER   = 80_000   # estimated tokens → auto-chunk above this
_CHUNK_SIZE      = 55_000   # tokens per map-phase chunk
_CHUNK_OVERLAP   = 400      # overlap tokens between chunks


def _estimate_tokens(text: str) -> int:
    return max(1, len(text) // _CHARS_PER_TOKEN)


def _split_text_chunks(text: str) -> list:
    """Split at paragraph boundaries into overlapping chunks."""
    cc = _CHUNK_SIZE    * _CHARS_PER_TOKEN
    oc = _CHUNK_OVERLAP * _CHARS_PER_TOKEN
    if len(text) <= cc:
        return [text]
    paras = text.split('\n\n')
    chunks, cur, cur_len = [], [], 0
    for p in paras:
        pl = len(p) + 2
        if cur_len + pl > cc and cur:
            full = '\n\n'.join(cur)
            chunks.append(full)
            tail = full[-oc:]
            cur, cur_len = ([tail] if tail else []), len(tail)
        cur.append(p)
        cur_len += pl
    if cur:
        chunks.append('\n\n'.join(cur))
    return chunks or [text]


async def _map_reduce(chunks: list, user_prompt: str, temperature: float,
                      max_tokens: int, key_row: dict, label: str) -> dict:
    """Map each chunk through the model, then reduce into a final answer."""
    n = min(len(chunks), 6)
    partial, pt, ct = [], 0, 0
    async with httpx.AsyncClient(timeout=300) as client:
        for i, chunk in enumerate(chunks[:n]):
            msgs = [
                {"role": "system",
                 "content": (f"You are analysing part {i+1} of {n} of a large document. "
                             "Extract only the key facts relevant to the user's request. Be concise.")},
                {"role": "user",
                 "content": f"Document part {i+1}:\n{chunk}\n\nRequest: {user_prompt}"},
            ]
            pay = _build_payload(msgs, min(max_tokens, 2048), temperature)
            r = await client.post(f"{settings.VLLM_URL}/v1/chat/completions", json=pay)
            if r.status_code == 200:
                rj = r.json()
                u  = rj.get("usage", {})
                pt += u.get("prompt_tokens", 0)
                ct += u.get("completion_tokens", 0)
                partial.append(
                    f"[Part {i+1}/{n}]\n{_strip_think_tags(_extract_text(rj['choices'][0]))}"
                )
    if not partial:
        raise HTTPException(502, "Map phase returned no results.")
    red_msgs = [
        {"role": "system",
         "content": ("Combine the following partial document analyses into one clear, "
                     "complete final answer. Remove redundancy. Preserve all unique facts.")},
        {"role": "user",
         "content": "\n\n".join(partial) + f"\n\nOriginal request: {user_prompt}"},
    ]
    pay2 = _build_payload(red_msgs, max_tokens, temperature)
    async with httpx.AsyncClient(timeout=300) as client:
        r2 = await client.post(f"{settings.VLLM_URL}/v1/chat/completions", json=pay2)
    if r2.status_code != 200:
        raise HTTPException(502, f"Reduce phase failed {r2.status_code}: {r2.text[:300]}")
    rj2 = r2.json()
    u2  = rj2.get("usage", {})
    pt += u2.get("prompt_tokens", 0)
    ct += u2.get("completion_tokens", 0)
    log_usage(key_row, label, pt, ct)
    ch = rj2["choices"][0]
    return {
        "text": _strip_think_tags(_extract_text(ch)),
        "finish_reason": ch.get("finish_reason", "stop"),
        "truncated": ch.get("finish_reason") == "length",
        "chunked": True,
        "chunks_processed": len(partial),
        "usage": {"input_tokens": pt, "output_tokens": ct, "total_tokens": pt + ct},
    }


# ──────────────────────────────────────────────────────────────────
#  Multi-format text extraction
#  Supported: .txt .md .csv .tsv .json .xml .html .yaml .toml
#             .xlsx .xls  (openpyxl)
#             .docx        (python-docx)
#             .pptx        (python-pptx)
# ──────────────────────────────────────────────────────────────────
_TEXT_EXTS = {
    ".txt", ".md", ".log", ".py", ".js", ".ts",
    ".json", ".xml", ".html", ".htm", ".css",
    ".yaml", ".yml", ".toml", ".ini", ".rst",
    ".csv", ".tsv",
    ".xlsx", ".xls",
    ".docx",
    ".pptx",
}


def _is_text_extractable(filename: str) -> bool:
    import os
    return os.path.splitext((filename or "").lower())[1] in _TEXT_EXTS


def _file_to_text(filename: str, content: bytes) -> str:
    """Extract plain text from office / document / plain-text files."""
    import os
    ext = os.path.splitext((filename or "").lower())[1]

    # ── Plain text family ──────────────────────────────────────────
    if ext in (".txt", ".md", ".log", ".py", ".js", ".ts",
               ".json", ".xml", ".html", ".htm", ".css",
               ".yaml", ".yml", ".toml", ".ini", ".rst"):
        for enc in ("utf-8", "utf-16", "latin-1"):
            try:
                return content.decode(enc)
            except Exception:
                pass
        return content.decode("utf-8", errors="replace")

    # ── CSV / TSV ──────────────────────────────────────────────────
    if ext in (".csv", ".tsv"):
        import csv, io
        sep = "\t" if ext == ".tsv" else ","
        reader = csv.reader(
            io.StringIO(content.decode("utf-8", errors="replace")), delimiter=sep
        )
        return "\n".join(" | ".join(row) for row in reader)

    # ── Excel ──────────────────────────────────────────────────────
    if ext in (".xlsx", ".xls"):
        try:
            import openpyxl, io as _io
            wb = openpyxl.load_workbook(
                _io.BytesIO(content), read_only=True, data_only=True
            )
            parts = []
            for name in wb.sheetnames:
                ws = wb[name]
                parts.append(f"=== Sheet: {name} ===")
                for row in ws.iter_rows(values_only=True):
                    parts.append(" | ".join("" if c is None else str(c) for c in row))
            return "\n".join(parts)
        except ImportError:
            raise HTTPException(
                500, "openpyxl not installed. Run: pip install openpyxl"
            )

    # ── Word ───────────────────────────────────────────────────────
    if ext == ".docx":
        try:
            import docx as _docx, io as _io
            doc = _docx.Document(_io.BytesIO(content))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            raise HTTPException(
                500, "python-docx not installed. Run: pip install python-docx"
            )

    # ── PowerPoint ────────────────────────────────────────────────
    if ext == ".pptx":
        try:
            import pptx as _pptx, io as _io
            prs = _pptx.Presentation(_io.BytesIO(content))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [
                    shape.text_frame.text
                    for shape in slide.shapes
                    if shape.has_text_frame and shape.text_frame.text.strip()
                ]
                slides.append(f"=== Slide {i} ===\n" + "\n".join(texts))
            return "\n\n".join(slides)
        except ImportError:
            raise HTTPException(
                500, "python-pptx not installed. Run: pip install python-pptx"
            )

    raise HTTPException(400, f"Cannot extract text from file type: {ext}")


# ──────────────────────────────────────────────────────────────────
#  /v1/generate (non-streaming)
# ──────────────────────────────────────────────────────────────────
@router.post("/v1/generate")
async def generate(req: GenerateRequest, key_row: dict = Depends(verify_key)):
    messages = []
    if req.system: messages.append({"role": "system", "content": req.system})
    messages.append({"role": "user", "content": req.prompt})
    payload = _build_payload(messages, req.max_tokens, req.temperature)
    try:
        async with httpx.AsyncClient(timeout=300) as client:
            resp = await client.post(f"{settings.VLLM_URL}/v1/chat/completions", json=payload)
        if resp.status_code != 200:
            log_usage(key_row, "/v1/generate", 0, 0, "error")
            raise HTTPException(status_code=502, detail=f"vLLM {resp.status_code}: {resp.text[:400]}")
        raw = resp.json()
        choice = raw["choices"][0]
        usage = raw.get("usage", {})
        log_usage(key_row, "/v1/generate", usage.get("prompt_tokens", 0), usage.get("completion_tokens", 0))
        return {
            "id": raw.get("id", ""),
            "model": settings.MODEL_DISPLAY_NAME,
            "text": _strip_think_tags(_extract_text(choice)),
            "finish_reason": choice.get("finish_reason", "stop"),
            "usage": {
                "input_tokens": usage.get("prompt_tokens", 0),
                "output_tokens": usage.get("completion_tokens", 0),
                "total_tokens": usage.get("total_tokens", 0),
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        log_usage(key_row, "/v1/generate", 0, 0, "error")
        raise HTTPException(status_code=500, detail=f"{type(e).__name__}: {str(e)[:300]}")


# ──────────────────────────────────────────────────────────────────
#  Streaming helper
# ──────────────────────────────────────────────────────────────────
async def _stream_vllm(payload, key_row, label, expose_thinking=False):
    """Stream vLLM response. Drops <think>...</think> blocks unless `expose_thinking=True`.
    When `expose_thinking=True`, thinking is sent as a separate `event:'thinking'` so the
    UI can show it in a collapsible panel without polluting the main answer."""
    payload = {**payload, "stream": True, "stream_options": {"include_usage": True}}
    final_usage = {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0}
    finish_reason = None
    think_filter = ThinkFilter()
    try:
        async with httpx.AsyncClient(timeout=600) as client:
            async with client.stream("POST", f"{settings.VLLM_URL}/v1/chat/completions",
                                     json=payload, headers={"Accept": "text/event-stream"}) as resp:
                if resp.status_code != 200:
                    err = (await resp.aread()).decode()[:400]
                    log_usage(key_row, label, 0, 0, "error")
                    yield f"data: {json.dumps({'event':'error','error':f'vLLM {resp.status_code}: {err}'})}\n\n"
                    return
                async for line in resp.aiter_lines():
                    if not line or not line.startswith("data:"): continue
                    data = line[5:].strip()
                    if data == "[DONE]":
                        # Flush any remaining filtered text
                        tail = think_filter.flush()
                        if tail:
                            yield f"data: {json.dumps({'event':'delta','text':tail})}\n\n"
                        log_usage(key_row, label, final_usage.get("prompt_tokens", 0), final_usage.get("completion_tokens", 0))
                        yield f"data: {json.dumps({'event':'done','finish_reason':finish_reason or 'stop','usage':{'input_tokens':final_usage.get('prompt_tokens',0),'output_tokens':final_usage.get('completion_tokens',0),'total_tokens':final_usage.get('total_tokens',0)},'model':settings.MODEL_DISPLAY_NAME})}\n\n"
                        return
                    try:
                        chunk = json.loads(data)
                    except json.JSONDecodeError:
                        continue
                    if chunk.get("usage"): final_usage = chunk["usage"]
                    choices = chunk.get("choices") or []
                    if not choices: continue
                    c = choices[0]
                    if c.get("finish_reason"): finish_reason = c["finish_reason"]
                    delta = c.get("delta") or {}

                    # Qwen3.6 emits separate `reasoning_content` field — surface as thinking
                    rc = delta.get("reasoning_content") or ""
                    if rc and expose_thinking:
                        yield f"data: {json.dumps({'event':'thinking','text':rc})}\n\n"

                    # Main content — strip inline <think> tags
                    raw = delta.get("content") or ""
                    if raw:
                        clean = think_filter.feed(raw)
                        if clean:
                            yield f"data: {json.dumps({'event':'delta','text':clean})}\n\n"
    except Exception as e:
        log_usage(key_row, label, 0, 0, "error")
        yield f"data: {json.dumps({'event':'error','error':f'{type(e).__name__}: {str(e)[:300]}'})}\n\n"


@router.post("/v1/generate/stream")
async def generate_stream(req: GenerateRequest, key_row: dict = Depends(verify_key)):
    messages = []
    if req.system: messages.append({"role": "system", "content": req.system})
    messages.append({"role": "user", "content": req.prompt})
    payload = _build_payload(messages, req.max_tokens, req.temperature)
    return StreamingResponse(_stream_vllm(payload, key_row, "/v1/generate/stream"),
                             media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache, no-transform", "Connection": "keep-alive", "X-Accel-Buffering": "no"})


@router.post("/v1/chat/completions")
async def chat(req: ChatRequest, key_row: dict = Depends(verify_key)):
    payload = {"model": settings.MODEL_NAME, "messages": req.messages, "max_tokens": req.max_tokens, "temperature": req.temperature}
    try:
        async with httpx.AsyncClient(timeout=300) as client:
            resp = await client.post(f"{settings.VLLM_URL}/v1/chat/completions", json=payload)
        result = resp.json()
        usage = result.get("usage", {})
        log_usage(key_row, "/v1/chat/completions", usage.get("prompt_tokens", 0), usage.get("completion_tokens", 0))
        return result
    except Exception as e:
        log_usage(key_row, "/v1/chat/completions", 0, 0, "error")
        raise HTTPException(status_code=500, detail=f"{type(e).__name__}: {str(e)[:300]}")


# ──────────────────────────────────────────────────────────────────
#  File → image data URLs
# ──────────────────────────────────────────────────────────────────
def _file_to_image_urls(filename, content, mime, dpi=200):
    """Convert PDF/image bytes → list of base64 data URLs.
    PDF: each page → PNG at given dpi (default 200 for accuracy).
    Image: passed through as-is.
    """
    name = (filename or "").lower()
    is_pdf = name.endswith(".pdf") or mime == "application/pdf"
    if is_pdf:
        try:
            import fitz
        except ImportError:
            raise HTTPException(500, "Run: pip install pymupdf")
        urls = []
        with fitz.open(stream=content, filetype="pdf") as doc:
            for page in doc:
                pix = page.get_pixmap(dpi=dpi)
                urls.append(f"data:image/png;base64,{base64.b64encode(pix.tobytes('png')).decode()}")
                if len(urls) >= 20:
                    break
        if not urls:
            raise HTTPException(400, "PDF had no pages.")
        return urls
    if not (mime or "").startswith("image/"):
        for ext, m in {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                       ".webp": "image/webp", ".gif": "image/gif", ".bmp": "image/bmp"}.items():
            if name.endswith(ext): mime = m; break
        else:
            mime = "image/jpeg"
    return [f"data:{mime};base64,{base64.b64encode(content).decode()}"]


# ──────────────────────────────────────────────────────────────────
#  /v1/vision/analyze (non-streaming)
# ──────────────────────────────────────────────────────────────────
@router.post("/v1/vision/analyze")
async def vision_analyze(
    file: UploadFile = File(...),
    prompt: str = Form(default="Describe this image."),
    max_tokens: int = Form(default=8192),
    temperature: float = Form(default=0.0),
    mode: Optional[str] = Form(default=None),
    key_row: dict = Depends(verify_key),
):
    try:
        content = await file.read()
        urls = _file_to_image_urls(file.filename or "", content, file.content_type or "")
        msg_content = [{"type": "image_url", "image_url": {"url": u}} for u in urls]
        msg_content.append({"type": "text", "text": prompt})
        messages = [{"role": "user", "content": msg_content}]
        payload = _build_payload(messages, max_tokens, temperature, mode=mode)
        async with httpx.AsyncClient(timeout=600) as client:
            resp = await client.post(f"{settings.VLLM_URL}/v1/chat/completions", json=payload)
        if resp.status_code != 200:
            log_usage(key_row, "/v1/vision/analyze", 0, 0, "error")
            raise HTTPException(502, f"vLLM {resp.status_code}: {resp.text[:400]}")
        raw = resp.json()
        choice = raw["choices"][0]
        text = _extract_text(choice)
        usage = raw.get("usage", {})
        log_usage(key_row, "/v1/vision/analyze", usage.get("prompt_tokens", 0), usage.get("completion_tokens", 0))
        text = _strip_think_tags(text)
        return {
            "id": raw.get("id", ""), "model": settings.MODEL_DISPLAY_NAME,
            "text": text, "result": text, "pages": len(urls), "mode": mode,
            "finish_reason": choice.get("finish_reason", "stop"),
            "usage": {
                "input_tokens": usage.get("prompt_tokens", 0),
                "output_tokens": usage.get("completion_tokens", 0),
                "total_tokens": usage.get("total_tokens", 0),
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        log_usage(key_row, "/v1/vision/analyze", 0, 0, "error")
        raise HTTPException(500, f"{type(e).__name__}: {str(e)[:300]}")


# ──────────────────────────────────────────────────────────────────
#  /v1/vision/analyze/stream — SSE + mode
# ──────────────────────────────────────────────────────────────────
@router.post("/v1/vision/analyze/stream")
async def vision_analyze_stream(
    file: UploadFile = File(...),
    prompt: str = Form(default="Describe this image."),
    max_tokens: int = Form(default=8192),
    temperature: float = Form(default=0.0),
    mode: Optional[str] = Form(default=None),
    key_row: dict = Depends(verify_key),
):
    content = await file.read()
    try:
        urls = _file_to_image_urls(file.filename or "", content, file.content_type or "")
    except HTTPException as e:
        async def err_gen():
            yield f"data: {json.dumps({'event':'error','error':e.detail})}\n\n"
        return StreamingResponse(err_gen(), media_type="text/event-stream")

    msg_content = [{"type": "image_url", "image_url": {"url": u}} for u in urls]
    msg_content.append({"type": "text", "text": prompt})
    messages = [{"role": "user", "content": msg_content}]
    payload = _build_payload(messages, max_tokens, temperature, mode=mode)

    async def stream_with_meta():
        yield f"data: {json.dumps({'event':'meta','pages':len(urls),'mode':mode or 'free'})}\n\n"
        async for chunk in _stream_vllm(payload, key_row, "/v1/vision/analyze/stream"):
            yield chunk

    return StreamingResponse(stream_with_meta(), media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache, no-transform", "Connection": "keep-alive", "X-Accel-Buffering": "no"})


# ──────────────────────────────────────────────────────────────────
#  /v1/modes — public list of active extraction modes (for frontend)
# ──────────────────────────────────────────────────────────────────
@router.get("/v1/modes")
def public_list_modes():
    rows = list_modes(include_inactive=False)
    return {
        "modes": [
            {"id": "free", "label": "Free-form", "icon": "✨",
             "description": "Standard prompt, no schema constraint",
             "default_prompt": ""}
        ] + [
            {
                "id": r["name"],
                "label": r["label"],
                "icon": r.get("icon") or "",
                "description": r.get("description") or "",
                "default_prompt": r.get("default_prompt") or "",
            }
            for r in rows
        ]
    }


# ══════════════════════════════════════════════════════════════════
#  UNIVERSAL ENDPOINTS — Gemini-style "one endpoint, many skills"
# ══════════════════════════════════════════════════════════════════

INTENT_KEYWORDS = {
    "invoice":  ["invoice", "gst", "bill of supply", "tax invoice"],
    "receipt":  ["receipt", "purchase", "bought", "store"],
    "id_card":  ["aadhaar", "pan card", "passport", "driving license", "id card"],
    "code":     ["function", "class", "algorithm", "python", "javascript", "typescript",
                 "go ", "rust", "java ", "sql", "regex", "code", "implement"],
    "translate":["translate", "translation", "in tamil", "in hindi", "in english",
                 "convert to "],
    "summarize":["summarize", "summary", "tldr", "tl;dr", "in short", "brief"],
}


def _detect_intent(prompt: str, has_file: bool) -> str:
    """Lightweight intent detection. Returns mode name or 'free'."""
    p = (prompt or "").lower()
    # If file attached, prioritise vision modes
    if has_file:
        for m in ("invoice", "receipt", "id_card"):
            if any(k in p for k in INTENT_KEYWORDS[m]):
                return m
        # Default vision intent if no specific keyword
        return "free"
    # Text-only — check keywords
    for mode, kws in INTENT_KEYWORDS.items():
        if mode in ("invoice", "receipt", "id_card"):
            continue  # only with file
        if any(k in p for k in kws):
            return mode
    return "free"


# ──────────────────────────────────────────────────────────────────
#  /v1/ask — UNIVERSAL POST endpoint
#  Accepts: prompt + optional file + optional mode (auto-detect if not given)
#  Returns: clean JSON response
# ──────────────────────────────────────────────────────────────────
@router.post("/v1/ask")
async def universal_ask(
    prompt: str = Form(...),
    file: Optional[UploadFile] = File(None),
    mode: Optional[str] = Form(None),
    max_tokens: int = Form(4096),
    temperature: float = Form(0.7),
    key_row: dict = Depends(verify_key),
):
    """One endpoint to rule them all.

    Examples:
      curl -X POST $BASE/v1/ask -H "X-API-Key: ak_..." \\
           -F "prompt=Hello, who are you?"

      curl -X POST $BASE/v1/ask -H "X-API-Key: ak_..." \\
           -F "prompt=Extract this invoice" -F "file=@invoice.pdf"

      curl -X POST $BASE/v1/ask -H "X-API-Key: ak_..." \\
           -F "prompt=Translate Hello to Tamil"
    """
    has_file = file is not None and file.filename
    auto_mode = mode or _detect_intent(prompt, has_file=has_file)

    try:
        if has_file:
            # Vision flow — use file + prompt + (auto) mode
            content = await file.read()
            urls = _file_to_image_urls(file.filename or "", content, file.content_type or "")
            msg_content = [{"type": "image_url", "image_url": {"url": u}} for u in urls]
            msg_content.append({"type": "text", "text": prompt})
            messages = [{"role": "user", "content": msg_content}]
            # If mode not vision-specific, drop it (use free-form)
            payload_mode = auto_mode if auto_mode in ("invoice", "receipt", "id_card") else None
            payload = _build_payload(messages, max_tokens, temperature, mode=payload_mode)
            label = "/v1/ask (vision)"
        else:
            # Text flow
            messages = [{"role": "user", "content": prompt}]
            payload = _build_payload(messages, max_tokens, temperature, mode=auto_mode if auto_mode != "free" else None)
            label = "/v1/ask (text)"

        async with httpx.AsyncClient(timeout=300) as client:
            resp = await client.post(f"{settings.VLLM_URL}/v1/chat/completions", json=payload)
        if resp.status_code != 200:
            log_usage(key_row, label, 0, 0, "error")
            raise HTTPException(502, f"vLLM {resp.status_code}: {resp.text[:300]}")
        raw = resp.json()
        choice = raw["choices"][0]
        text = _strip_think_tags(_extract_text(choice))
        usage = raw.get("usage", {})
        log_usage(key_row, label, usage.get("prompt_tokens", 0), usage.get("completion_tokens", 0))
        return {
            "id": raw.get("id", ""),
            "model": settings.MODEL_DISPLAY_NAME,
            "text": text,
            "detected_intent": auto_mode,
            "explicit_mode": mode,
            "had_file": has_file,
            "finish_reason": choice.get("finish_reason", "stop"),
            "usage": {
                "input_tokens": usage.get("prompt_tokens", 0),
                "output_tokens": usage.get("completion_tokens", 0),
                "total_tokens": usage.get("total_tokens", 0),
            },
        }
    except HTTPException:
        raise
    except Exception as e:
        log_usage(key_row, "/v1/ask", 0, 0, "error")
        raise HTTPException(500, f"{type(e).__name__}: {str(e)[:300]}")


# ──────────────────────────────────────────────────────────────────
#  /v1/quick — GET endpoint for casual browser testing ONLY
#  ⚠ NOT for production — prompts visible in URL/logs
#  Limit: 500 char prompt, no file uploads possible via GET
# ──────────────────────────────────────────────────────────────────
@router.get("/v1/quick")
async def quick_ask(
    q: str,
    mode: Optional[str] = None,
    max_tokens: int = 512,
    key_row: dict = Depends(verify_key),
):
    """Quick GET endpoint — paste in browser to test.

    Example URL (in browser address bar — replace YOUR_KEY):
      https://wo50dppqmt72bl-1111.proxy.runpod.net/v1/quick?q=Hello&max_tokens=100

    Headers required:
      X-API-Key: ak_...   (use a browser extension to set, OR use POST /v1/ask)

    Limits:
      - Max 500 chars in `q`
      - Max 512 tokens in response
      - No file uploads (impossible via GET)
      - Logged in plain URL — DO NOT use for sensitive prompts
    """
    if len(q) > 500:
        raise HTTPException(400, "Prompt too long for GET (max 500 chars). Use POST /v1/ask for longer prompts.")
    if max_tokens > 1024:
        raise HTTPException(400, "GET endpoint capped at max_tokens=1024. Use POST /v1/ask for more.")

    auto_mode = mode or _detect_intent(q, has_file=False)
    messages = [{"role": "user", "content": q}]
    payload = _build_payload(messages, max_tokens, 0.7,
                             mode=auto_mode if auto_mode != "free" else None)
    try:
        async with httpx.AsyncClient(timeout=120) as client:
            resp = await client.post(f"{settings.VLLM_URL}/v1/chat/completions", json=payload)
        if resp.status_code != 200:
            log_usage(key_row, "/v1/quick", 0, 0, "error")
            raise HTTPException(502, f"vLLM {resp.status_code}: {resp.text[:200]}")
        raw = resp.json()
        choice = raw["choices"][0]
        text = _strip_think_tags(_extract_text(choice))
        usage = raw.get("usage", {})
        log_usage(key_row, "/v1/quick", usage.get("prompt_tokens", 0), usage.get("completion_tokens", 0))
        return {
            "answer": text,
            "detected_intent": auto_mode,
            "tokens": usage.get("total_tokens", 0),
            "warning": "GET endpoint — prompts logged in URL. Use POST /v1/ask for production.",
        }
    except HTTPException:
        raise
    except Exception as e:
        log_usage(key_row, "/v1/quick", 0, 0, "error")
        raise HTTPException(500, f"{type(e).__name__}: {str(e)[:200]}")
